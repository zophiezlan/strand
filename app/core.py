"""
Core hair rendering and bytes-based injectors for Strand.

Shared by the FastAPI web service (`app.main`) and the CLI (`app.cli`). The
public API is:

    InjectOptions       — request-level knobs
    strand_bytes(...)   — dispatch on filename / content type
    strand_zip_bytes(...) — walk a zip, strand each supported entry
    inject_image_bytes(...), inject_pdf_bytes(...), inject_pptx_bytes(...)

Everything operates on bytes / BytesIO. Nothing in this module touches the disk.
"""

from __future__ import annotations

import io
import math
import random
from dataclasses import dataclass, field

from PIL import Image, ImageDraw


# Supersample at 4x and downsample for smooth strand edges.
_SCALE = 4
_PAD = 12


def _cubic_bezier(t, p0, p1, p2, p3):
    mt = 1 - t
    x = mt**3 * p0[0] + 3 * mt**2 * t * p1[0] + 3 * mt * t**2 * p2[0] + t**3 * p3[0]
    y = mt**3 * p0[1] + 3 * mt**2 * t * p1[1] + 3 * mt * t**2 * p2[1] + t**3 * p3[1]
    return (x, y)


def _new_canvas(base_width, base_height):
    cw = (base_width + 2 * _PAD) * _SCALE
    ch = (base_height + 2 * _PAD) * _SCALE
    img = Image.new("RGBA", (cw, ch), (0, 0, 0, 0))
    return img, ImageDraw.Draw(img), cw, ch


def _to_canvas(x, y):
    return ((_PAD + x) * _SCALE, (_PAD + y) * _SCALE)


PALETTES = {
    "dark":    ((10, 6, 4),     (45, 32, 28)),
    "brown":   ((45, 30, 18),   (95, 65, 38)),
    "blonde":  ((130, 100, 55), (210, 180, 130)),
    "grey":    ((90, 88, 85),   (170, 168, 165)),
    "white":   ((200, 198, 195),(235, 233, 230)),
    "red":     ((90, 38, 18),   (165, 80, 42)),
}
PALETTE_NAMES = set(PALETTES) | {"mixed"}

# Per-palette stroke width scaling. A real hair is ~80 µm wide — sub-pixel at
# typical DPI — so any visible thickness is already a perceptual lie. White
# hair on white paper has near-zero contrast so the eye doesn't measure its
# width; dark hair has high contrast and every pixel of thickness reads as
# "drawn line" instead of "scanned hair". This table damps darker palettes
# down toward sub-pixel territory so they feel like strands, not strokes.
PALETTE_WIDTH_SCALE: dict[str, float] = {
    "dark":   0.55,
    "brown":  0.65,
    "red":    0.65,
    "blonde": 0.75,
    "grey":   0.75,
    "white":  1.00,
}


def _palette_width_scale(palette: str) -> float:
    return PALETTE_WIDTH_SCALE.get(palette, 1.0)


def _random_hair_color(rng, palette="dark"):
    if palette == "mixed":
        palette = rng.choice(list(PALETTES))
    lo, hi = PALETTES[palette]
    return (rng.randint(lo[0], hi[0]),
            rng.randint(lo[1], hi[1]),
            rng.randint(lo[2], hi[2]))


# Probability that a strand gets a follicle bulb at one end. Real shed hair
# often has a tiny darker "club" root at the end that fell out — it's the
# single most diagnostic "yep, real hair" detail. Eyelashes get a lower
# chance (eyelashes are rarely shed root-and-all in everyday context); broken
# fragments never get one (a broken-mid-shaft end isn't a root).
_FOLLICLE_CHANCE = 0.18
_FOLLICLE_CHANCE_EYELASH = 0.10


def _draw_follicle(draw, position, base_color, rng, width_scale=1.0):
    """Composite a small dark bulb at `position` to suggest a hair follicle root.

    `position` is in canvas coordinates (already supersampled). The bulb is a
    short oval, ~2× the local strand width, in a colour darker than the
    strand. Scaled by `width_scale` so dark/thin palettes don't get a chunky
    contrast-grabbing dot.
    """
    # Darken the strand colour by 50% — visible without being a black blob.
    r = max(0, int(base_color[0] * 0.45))
    g = max(0, int(base_color[1] * 0.45))
    b = max(0, int(base_color[2] * 0.45))
    a = rng.randint(220, 245)

    # Bulb radius in final-image px scaled to the supersample canvas.
    radius = (rng.uniform(1.8, 2.6) * width_scale) * _SCALE
    # Slight asymmetry so it reads as an oval, not a perfect circle.
    rx = radius * rng.uniform(1.0, 1.25)
    ry = radius * rng.uniform(0.8, 1.0)
    cx, cy = position
    draw.ellipse(
        [cx - rx, cy - ry, cx + rx, cy + ry],
        fill=(r, g, b, a),
    )


def _draw_strand_polyline(draw, points, base_color, rng,
                          width_start=1.9, width_end=1.0, width_scale=1.0):
    """Draw a strand following an arbitrary polyline of canvas-coord points.

    Width tapers from `width_start` to `width_end` along the polyline with a
    small per-segment jitter; colour gets per-pixel jitter so the strand reads
    as a real photographed object rather than a vector stroke. Shared by the
    bezier-based morphologies and the kink (sine-modulated polyline) one.
    """
    width_start *= width_scale
    width_end *= width_scale
    n = len(points) - 1
    if n <= 0:
        return
    for i in range(n):
        t = i / n
        w = (width_start - t * (width_start - width_end) + rng.uniform(-0.15, 0.15)) * _SCALE
        w = max(_SCALE * 0.5, w)

        r = max(0, min(255, base_color[0] + rng.randint(-8, 8)))
        g = max(0, min(255, base_color[1] + rng.randint(-6, 6)))
        b = max(0, min(255, base_color[2] + rng.randint(-5, 5)))
        a = rng.randint(220, 250)

        draw.line([points[i], points[i + 1]], fill=(r, g, b, a), width=int(round(w)))


def _draw_bezier_segment(draw, p0, p1, p2, p3, base_color, rng,
                         n_samples=220, width_start=1.9, width_end=1.0,
                         width_scale=1.0):
    points = [_cubic_bezier(i / (n_samples - 1), p0, p1, p2, p3) for i in range(n_samples)]
    _draw_strand_polyline(draw, points, base_color, rng,
                          width_start=width_start, width_end=width_end,
                          width_scale=width_scale)


def _finalize(img, cw, ch):
    return img.resize((cw // _SCALE, ch // _SCALE), Image.LANCZOS)


def _generate_curve_hair(rng, palette="dark", base_width=400, base_height=120):
    img, draw, cw, ch = _new_canvas(base_width, base_height)

    p0 = _to_canvas(
        rng.uniform(0, base_width * 0.05),
        base_height * 0.5 + rng.uniform(-base_height * 0.2, base_height * 0.2),
    )
    p3 = _to_canvas(
        base_width * 0.95 + rng.uniform(-base_width * 0.05, base_width * 0.03),
        base_height * 0.5 + rng.uniform(-base_height * 0.3, base_height * 0.3),
    )
    p1 = _to_canvas(
        base_width * 0.3 + rng.uniform(-base_width * 0.1, base_width * 0.1),
        base_height * 0.5 + rng.uniform(-base_height * 0.7, base_height * 0.7),
    )
    p2 = _to_canvas(
        base_width * 0.7 + rng.uniform(-base_width * 0.1, base_width * 0.1),
        base_height * 0.5 + rng.uniform(-base_height * 0.7, base_height * 0.7),
    )

    base_color = _random_hair_color(rng, palette)
    ws = _palette_width_scale(palette)
    _draw_bezier_segment(draw, p0, p1, p2, p3, base_color, rng,
                         n_samples=220, width_start=1.9, width_end=1.0,
                         width_scale=ws)
    if rng.random() < _FOLLICLE_CHANCE:
        _draw_follicle(draw, rng.choice([p0, p3]), base_color, rng, width_scale=ws)
    return _finalize(img, cw, ch)


def _generate_loop_hair(rng, palette="dark", base_width=400, base_height=240):
    img, draw, cw, ch = _new_canvas(base_width, base_height)

    loop_cx = base_width * rng.uniform(0.25, 0.38)
    loop_cy = base_height * 0.5
    loop_r = min(base_width, base_height) * rng.uniform(0.22, 0.32)

    theta_close = rng.uniform(-0.35, 0.35)
    theta_open = theta_close + rng.uniform(0.18, 0.45)

    a_local = (loop_cx + loop_r * math.cos(theta_close),
               loop_cy + loop_r * math.sin(theta_close))
    b_local = (loop_cx + loop_r * math.cos(theta_open),
               loop_cy + loop_r * math.sin(theta_open))

    tan_a = (-math.sin(theta_close), math.cos(theta_close))
    tan_b = (math.sin(theta_open), -math.cos(theta_open))

    # ~1.7r gives a near-circular Bezier loop
    k = loop_r * 1.7
    p1_local = (a_local[0] + k * tan_a[0], a_local[1] + k * tan_a[1])
    p2_local = (b_local[0] + k * tan_b[0], b_local[1] + k * tan_b[1])

    base_color = _random_hair_color(rng, palette)
    ws = _palette_width_scale(palette)
    _draw_bezier_segment(
        draw,
        _to_canvas(*a_local), _to_canvas(*p1_local),
        _to_canvas(*p2_local), _to_canvas(*b_local),
        base_color, rng,
        n_samples=200, width_start=1.4, width_end=1.4,
        width_scale=ws,
    )

    exit_dx = b_local[0] - p2_local[0]
    exit_dy = b_local[1] - p2_local[1]
    mag = math.hypot(exit_dx, exit_dy) or 1.0
    tail_dir = (exit_dx / mag, exit_dy / mag)

    tail_len = base_width * rng.uniform(0.40, 0.55)
    t_end_local = (b_local[0] + tail_len * tail_dir[0] + rng.uniform(-15, 15),
                   b_local[1] + tail_len * tail_dir[1] + rng.uniform(-25, 25))
    t_p1_local = (b_local[0] + tail_len * 0.3 * tail_dir[0] + rng.uniform(-10, 10),
                  b_local[1] + tail_len * 0.3 * tail_dir[1] + rng.uniform(-12, 12))
    t_p2_local = (b_local[0] + tail_len * 0.7 * tail_dir[0] + rng.uniform(-15, 15),
                  b_local[1] + tail_len * 0.7 * tail_dir[1] + rng.uniform(-18, 18))

    _draw_bezier_segment(
        draw,
        _to_canvas(*b_local), _to_canvas(*t_p1_local),
        _to_canvas(*t_p2_local), _to_canvas(*t_end_local),
        base_color, rng,
        n_samples=140, width_start=1.4, width_end=0.9,
        width_scale=ws,
    )
    if rng.random() < _FOLLICLE_CHANCE:
        # Bulb sits at the trailing tail end — that's the "torn off the head"
        # end of a shed strand; the loop is mid-shaft.
        _draw_follicle(draw, _to_canvas(*t_end_local), base_color, rng, width_scale=ws)

    return _finalize(img, cw, ch)


def _generate_eyelash_hair(rng: random.Random, palette: str = "dark",
                           base_width: int = 400, base_height: int = 120) -> Image.Image:
    """A short, tightly curved hair — eyelash-like.

    Drawn in a small region of a full-sized canvas, so when `_place` scales the
    final image, an eyelash naturally ends up smaller on the page than a curve
    or loop. (The transparent surround does the work.)
    """
    img, draw, cw, ch = _new_canvas(base_width, base_height)

    cx = base_width * 0.5
    cy = base_height * 0.5

    span = rng.uniform(28, 50)
    arc = rng.uniform(18, 32) * rng.choice([-1, 1])

    p0 = _to_canvas(cx - span / 2, cy + rng.uniform(-2, 2))
    p3 = _to_canvas(cx + span / 2, cy + rng.uniform(-2, 2))
    p1 = _to_canvas(cx - span * 0.18, cy + arc)
    p2 = _to_canvas(cx + span * 0.18, cy + arc * 0.85)

    base_color = _random_hair_color(rng, palette)
    ws = _palette_width_scale(palette)
    _draw_bezier_segment(draw, p0, p1, p2, p3, base_color, rng,
                         n_samples=120, width_start=1.9, width_end=1.0,
                         width_scale=ws)
    if rng.random() < _FOLLICLE_CHANCE_EYELASH:
        _draw_follicle(draw, rng.choice([p0, p3]), base_color, rng, width_scale=ws)
    return _finalize(img, cw, ch)


def _generate_kink_hair(rng: random.Random, palette: str = "dark",
                        base_width: int = 400, base_height: int = 120) -> Image.Image:
    """A tightly coiled hair — pubic / curly body-hair vibe.

    Real curly hair has *three* scales of shape stacked:
      • a wandering underlying path (the strand itself isn't laid out straight
        on the surface — it wanders chaotically), modelled as a cubic bezier
      • a slow curl along that path (the big "this is curly hair" arcs)
      • a fast crimp riding on top (the tight zigzag texture)

    The two sine modulations are offset perpendicular to the bezier's
    *tangent* at each sample, so the kinks ride correctly on the curving
    underlying path rather than along a fixed axis. Both modulations are
    bell-enveloped at the ends so the strand doesn't appear cut off.
    """
    img, draw, cw, ch = _new_canvas(base_width, base_height)

    cx_start = base_width * rng.uniform(0.08, 0.18)
    cx_end = base_width * rng.uniform(0.78, 0.92)
    cy_mid = base_height * 0.5
    cy_start = cy_mid + rng.uniform(-base_height * 0.20, base_height * 0.20)
    cy_end = cy_mid + rng.uniform(-base_height * 0.20, base_height * 0.20)

    # Wacky primary path: cubic bezier with control points that can swing
    # ±35% of base_height off the midline, so the underlying strand wanders.
    p0 = (cx_start, cy_start)
    p3 = (cx_end, cy_end)
    span_x = cx_end - cx_start
    p1 = (cx_start + span_x * rng.uniform(0.20, 0.40),
          cy_mid + rng.uniform(-base_height * 0.35, base_height * 0.35))
    p2 = (cx_start + span_x * rng.uniform(0.60, 0.80),
          cy_mid + rng.uniform(-base_height * 0.35, base_height * 0.35))

    # Approximate arc length of the bezier so sine amplitudes scale to the
    # actual path length, not the straight-line distance between endpoints.
    coarse = [_cubic_bezier(k / 24, p0, p1, p2, p3) for k in range(25)]
    arc_len = sum(math.hypot(coarse[k + 1][0] - coarse[k][0],
                             coarse[k + 1][1] - coarse[k][1])
                  for k in range(24)) or 1.0

    # Slow curl along the wandering path.
    slow_cycles = rng.uniform(0.6, 1.6)
    slow_amp = arc_len * rng.uniform(0.10, 0.18)
    slow_phase = rng.uniform(0.0, math.tau)

    # Fast crimp on top.
    fast_cycles = rng.uniform(5.0, 9.0)
    fast_amp = arc_len * rng.uniform(0.035, 0.075)
    fast_phase = rng.uniform(0.0, math.tau)

    n_samples = 280
    points = []
    for i in range(n_samples):
        t = i / (n_samples - 1)
        # Point on bezier
        bx, by = _cubic_bezier(t, p0, p1, p2, p3)
        # Tangent (first derivative of cubic bezier)
        mt = 1 - t
        tx = (3 * mt * mt * (p1[0] - p0[0]) + 6 * mt * t * (p2[0] - p1[0])
              + 3 * t * t * (p3[0] - p2[0]))
        ty = (3 * mt * mt * (p1[1] - p0[1]) + 6 * mt * t * (p2[1] - p1[1])
              + 3 * t * t * (p3[1] - p2[1]))
        tlen = math.hypot(tx, ty) or 1.0
        # Unit perpendicular to the tangent — the kinks ride sideways on the path
        px_dir = -ty / tlen
        py_dir = tx / tlen
        # Bell envelope tapers amplitude at the ends.
        env = math.sin(math.pi * t) ** 0.5
        offset = (
            slow_amp * env * math.sin(t * slow_cycles * math.tau + slow_phase)
            + fast_amp * env * math.sin(t * fast_cycles * math.tau + fast_phase)
        )
        points.append(_to_canvas(bx + px_dir * offset, by + py_dir * offset))

    base_color = _random_hair_color(rng, palette)
    ws = _palette_width_scale(palette)
    _draw_strand_polyline(draw, points, base_color, rng,
                          width_start=1.6, width_end=0.9,
                          width_scale=ws)
    if rng.random() < _FOLLICLE_CHANCE:
        _draw_follicle(draw, rng.choice([points[0], points[-1]]),
                       base_color, rng, width_scale=ws)
    return _finalize(img, cw, ch)


def _generate_fragment_hair(rng: random.Random, palette: str = "dark",
                            base_width: int = 400, base_height: int = 120) -> Image.Image:
    """A short broken piece of hair — fragmentary, less curved than an eyelash."""
    img, draw, cw, ch = _new_canvas(base_width, base_height)

    cx = base_width * 0.5
    cy = base_height * 0.5

    span = rng.uniform(45, 90)
    bend = rng.uniform(-12, 12)

    p0 = _to_canvas(cx - span / 2 + rng.uniform(-3, 3), cy + rng.uniform(-3, 3))
    p3 = _to_canvas(cx + span / 2 + rng.uniform(-3, 3), cy + rng.uniform(-5, 5))
    p1 = _to_canvas(cx - span * 0.2, cy + bend)
    p2 = _to_canvas(cx + span * 0.2, cy + bend * 0.6 + rng.uniform(-4, 4))

    base_color = _random_hair_color(rng, palette)
    # Thinner than a full strand — fragments are wispy.
    _draw_bezier_segment(draw, p0, p1, p2, p3, base_color, rng,
                         n_samples=110, width_start=1.4, width_end=0.7,
                         width_scale=_palette_width_scale(palette))
    return _finalize(img, cw, ch)


# Public dispatch table — used by /api/sample to render a specific morphology
# and by `generate_hair` for weighted random selection.
MORPHOLOGIES: dict[str, "callable"] = {
    "curve":    _generate_curve_hair,
    "loop":     _generate_loop_hair,
    "eyelash":  _generate_eyelash_hair,
    "fragment": _generate_fragment_hair,
    "kink":     _generate_kink_hair,
}


# Real-world length ranges per morphology, in centimetres. Used by `_place`
# together with the substrate's units-per-cm to size each hair to a physically
# plausible length — a head-hair strand is 3-12 cm long whether it lands on a
# passport photo or an A4 scan, an eyelash is always tiny.
MORPHOLOGY_LENGTH_CM: dict[str, tuple[float, float]] = {
    "curve":    (3.0, 12.0),
    "loop":     (5.0, 18.0),
    "eyelash":  (0.7, 1.4),
    "fragment": (1.5, 4.0),
    "kink":     (1.5, 4.0),     # pubic / tightly coiled body hair
}

# Substrate unit conversions per centimetre.
_POINTS_PER_CM = 72.0 / 2.54        # PDFs use points (72 pt = 1 in)
_EMU_PER_CM = 914400.0 / 2.54       # pptx uses EMU (914400 EMU = 1 in)
_DEFAULT_DPI = 96                   # screen-typical fallback for images
# Soft cap to keep absurdly high DPI metadata (e.g. 1200 dpi on a phone
# screenshot) from producing tiny hairs that don't read as hairs.
_MAX_IMAGE_DPI = 300


def _image_pixels_per_cm(img: Image.Image) -> float:
    """Read DPI from PIL image metadata; fall back to 96 and cap at 300."""
    dpi_info = img.info.get("dpi")
    if dpi_info and isinstance(dpi_info, tuple) and dpi_info[0]:
        dpi = float(dpi_info[0])
    else:
        dpi = float(_DEFAULT_DPI)
    dpi = min(dpi, _MAX_IMAGE_DPI)
    return dpi / 2.54


def generate_hair_with_morphology(
    rng: random.Random,
    palette: str = "dark",
    loop_chance: float = 0.15,
    eyelash_chance: float = 0.08,
    fragment_chance: float = 0.08,
    kink_chance: float = 0.06,
) -> tuple[Image.Image, str, str]:
    """Pick a morphology by weighted random choice; return (image, morphology, palette).

    For "mixed" palette, a concrete palette is resolved once per hair (so one
    strand stays in a single colour family) and the resolved name is returned
    so callers can build a colour breakdown.
    """
    r = rng.random()
    # Morphology pick happens first so the RNG sequence is unchanged when
    # palette isn't "mixed" (preserves the seed-determinism guarantee).
    if r < loop_chance:
        morph = "loop"
        renderer = _generate_loop_hair
    elif r < loop_chance + eyelash_chance:
        morph = "eyelash"
        renderer = _generate_eyelash_hair
    elif r < loop_chance + eyelash_chance + fragment_chance:
        morph = "fragment"
        renderer = _generate_fragment_hair
    elif r < loop_chance + eyelash_chance + fragment_chance + kink_chance:
        morph = "kink"
        renderer = _generate_kink_hair
    else:
        morph = "curve"
        renderer = _generate_curve_hair

    # Resolve "mixed" → concrete palette here (rather than inside
    # _random_hair_color) so we can report which colour family was drawn.
    resolved = palette
    if resolved == "mixed":
        resolved = rng.choice(list(PALETTES))

    return renderer(rng, palette=resolved), morph, resolved


def generate_hair(
    rng: random.Random,
    palette: str = "dark",
    loop_chance: float = 0.15,
    eyelash_chance: float = 0.08,
    fragment_chance: float = 0.08,
    kink_chance: float = 0.06,
) -> Image.Image:
    """Pick a morphology by weighted random choice; default is the curved strand."""
    img, _, _ = generate_hair_with_morphology(
        rng, palette=palette,
        loop_chance=loop_chance, eyelash_chance=eyelash_chance,
        fragment_chance=fragment_chance, kink_chance=kink_chance,
    )
    return img


def _empty_stats() -> dict:
    """Fresh stats dict — every counter the injectors fill in.

    `palettes` and `morphologies` are dynamic maps of name → count.
    `clusters` counts primary hairs that grew at least one buddy.
    `hair_lengths_cm` is appended to once per hair drawn (primary + buddy).
    `content_hits` counts primary hairs that landed via the content-region bias.
    `substrate` is set once per file (None inside the zip aggregate, since a
    zip mixes multiple substrates).
    """
    return {
        "hairs": 0,
        "morphologies": {"curve": 0, "loop": 0, "eyelash": 0, "fragment": 0, "kink": 0},
        "palettes": {},
        "pages_touched": 0,
        "clusters": 0,
        "hair_lengths_cm": [],
        "content_hits": 0,
        "substrate": None,
    }


def _bump(d: dict, key: str) -> None:
    d[key] = d.get(key, 0) + 1


def _merge_stats(into: dict, src: dict) -> None:
    """Aggregate `src` stats into `into` in place. Used by the zip handler."""
    into["hairs"] += src.get("hairs", 0)
    into["pages_touched"] += src.get("pages_touched", 0)
    into["clusters"] += src.get("clusters", 0)
    into["content_hits"] += src.get("content_hits", 0)
    into["hair_lengths_cm"].extend(src.get("hair_lengths_cm", []))
    for m, n in src.get("morphologies", {}).items():
        into["morphologies"][m] = into["morphologies"].get(m, 0) + n
    for p, n in src.get("palettes", {}).items():
        into["palettes"][p] = into["palettes"].get(p, 0) + n
    # `substrate` deliberately stays None on the aggregate — a zip has many.


def _rotate_hair(hair: Image.Image, rng: random.Random) -> Image.Image:
    angle = rng.uniform(0, 360)
    return hair.rotate(angle, expand=True, resample=Image.BICUBIC)


# ---------------------------------------------------------------------------
# Content region detection — bias placement toward areas with real content.
# ---------------------------------------------------------------------------

def _image_content_regions(img: Image.Image, grid: int = 10, threshold: float = 3.0):
    from PIL import ImageFilter, ImageStat
    gray = img.convert("L")
    edges = gray.filter(ImageFilter.FIND_EDGES)
    w, h = edges.size
    cw = max(1, w // grid)
    ch = max(1, h // grid)
    regions = []
    for cy in range(grid):
        for cx in range(grid):
            x0 = cx * cw
            y0 = cy * ch
            x1 = x0 + cw
            y1 = y0 + ch
            cell = edges.crop((x0, y0, x1, y1))
            mean = ImageStat.Stat(cell).mean[0]
            if mean >= threshold:
                regions.append((x0, y0, x1, y1, mean))
    return regions


def _pdf_content_regions(page):
    regions = []
    try:
        blocks = page.get_text("blocks")
    except Exception:
        return regions
    for b in blocks:
        if len(b) < 4:
            continue
        x0, y0, x1, y1 = b[:4]
        area = max(1.0, (x1 - x0) * (y1 - y0))
        regions.append((x0, y0, x1, y1, area))
    return regions


def _pptx_content_regions(slide):
    regions = []
    for shape in slide.shapes:
        if (shape.left is None or shape.top is None or
                shape.width is None or shape.height is None):
            continue
        x0 = shape.left
        y0 = shape.top
        x1 = x0 + shape.width
        y1 = y0 + shape.height
        area = max(1.0, shape.width * shape.height)
        regions.append((x0, y0, x1, y1, area))
    return regions


def _place(rng, container_w, container_h, hair_w, hair_h,
           length_range_cm, units_per_cm,
           regions=None, content_bias: float = 0.5):
    """Pick a position + size for one hair, sized to a real-world length.

    `length_range_cm` is the (min, max) physical length in centimetres for the
    morphology; `units_per_cm` converts cm into the substrate's native units
    (pixels for images, points for PDFs, EMU for pptx). The hair's longest
    edge is sized to a random length in that range, then clamped to <= 90% of
    the container's short side so a hair on a thumbnail can't dominate it.
    """
    target_length_cm = rng.uniform(*length_range_cm)
    target_long = target_length_cm * units_per_cm
    short_side = min(container_w, container_h)
    # Don't let a hair eat the whole substrate (thumbnails, narrow strips).
    target_long = min(target_long, short_side * 0.9)
    scale = target_long / max(hair_w, hair_h)
    dw = hair_w * scale
    dh = hair_h * scale

    use_region = bool(regions) and rng.random() < content_bias

    if use_region:
        weights = [r[4] for r in regions]
        region = rng.choices(regions, weights=weights, k=1)[0]
        rx0, ry0, rx1, ry1 = region[:4]
        rcx = (rx0 + rx1) / 2
        rcy = (ry0 + ry1) / 2
        jitter_x = max((rx1 - rx0) * 0.6, dw * 0.2)
        jitter_y = max((ry1 - ry0) * 0.6, dh * 0.2)
        x = rcx - dw / 2 + rng.uniform(-jitter_x, jitter_x)
        y = rcy - dh / 2 + rng.uniform(-jitter_y, jitter_y)
    else:
        max_x = max(0.0, container_w - dw)
        max_y = max(0.0, container_h - dh)
        x = rng.triangular(0, max_x, max_x / 2) if max_x > 0 else 0
        y = rng.triangular(0, max_y, max_y / 2) if max_y > 0 else 0

    x = max(0, min(container_w - dw, x))
    y = max(0, min(container_h - dh, y))
    return x, y, dw, dh, use_region


def _morph_kwargs(opts) -> dict:
    """Bundle morphology weights for a `generate_hair` call."""
    return dict(
        palette=opts.palette,
        loop_chance=opts.loop_chance,
        eyelash_chance=opts.eyelash_chance,
        fragment_chance=opts.fragment_chance,
        kink_chance=opts.kink_chance,
    )


def _buddy_offsets(rng, dw, dh, container_w, container_h, base_x, base_y,
                   cluster_chance, max_buddies=2):
    """
    Roll up to `max_buddies` times for a buddy hair near (base_x, base_y).

    Each roll succeeds with probability `cluster_chance`. A successful roll
    yields an (x, y) offset within ~0.6× hair size of the seed position,
    clamped to the container bounds. Lets real-world clumping show up
    without requiring every hair to clump.
    """
    out = []
    for _ in range(max_buddies):
        if rng.random() >= cluster_chance:
            break
        ox = base_x + rng.uniform(-dw * 0.6, dw * 0.6)
        oy = base_y + rng.uniform(-dh * 0.6, dh * 0.6)
        ox = max(0, min(container_w - dw, ox))
        oy = max(0, min(container_h - dh, oy))
        out.append((ox, oy))
    return out


# ---------------------------------------------------------------------------
# Public API
# ---------------------------------------------------------------------------

IMAGE_SUFFIXES = {".png", ".jpg", ".jpeg", ".gif", ".bmp", ".webp", ".tif", ".tiff", ".ico"}
PDF_SUFFIXES = {".pdf"}
PPTX_SUFFIXES = {".pptx"}
DOCX_SUFFIXES = {".docx"}
XLSX_SUFFIXES = {".xlsx"}
ZIP_SUFFIXES = {".zip"}
SUPPORTED_SUFFIXES = (
    IMAGE_SUFFIXES | PDF_SUFFIXES | PPTX_SUFFIXES | DOCX_SUFFIXES | XLSX_SUFFIXES
)
SUPPORTED_INCLUDING_ZIP = SUPPORTED_SUFFIXES | ZIP_SUFFIXES


# Ordered intensity tiers. Past "heavy" the labels are the joke — every page
# gets multiple hairs, every image gets a fistful.
INTENSITY_TIERS: dict[str, dict] = {
    "subtle":     {"label": "Subtle",     "image_count": 1,  "rate": 0.25, "hairs_per_page": 1},
    "normal":     {"label": "Normal",     "image_count": 2,  "rate": 0.50, "hairs_per_page": 1},
    "heavy":      {"label": "Heavy",      "image_count": 4,  "rate": 0.85, "hairs_per_page": 2},
    "hirsute":    {"label": "Hirsute",    "image_count": 10, "rate": 1.0,  "hairs_per_page": 4},
    "werewolf":   {"label": "Werewolf",   "image_count": 25, "rate": 1.0,  "hairs_per_page": 9},
    "cousin-itt": {"label": "Cousin Itt", "image_count": 60, "rate": 1.0,  "hairs_per_page": 22},
}
INTENSITY_ORDER: list[str] = list(INTENSITY_TIERS.keys())
# Default tuned by real-world testing: a single light hair best mimics the
# photocopier-glass / laminator-pocket artefact, which reads most convincingly
# as a "real" stray hair. Heavier tiers exist for play, not for first impressions.
DEFAULT_INTENSITY = "subtle"
DEFAULT_PALETTE = "white"


@dataclass
class InjectOptions:
    """Request-level knobs for a single strand call."""
    # Per-page/slide probability for PDFs and pptx. For images, see image_count.
    rate: float = 0.5
    # Hairs per *selected* PDF page or pptx slide.
    hairs_per_page: int = 1
    # Hairs to overlay on a single image. PDFs/pptx use rate + hairs_per_page.
    image_count: int = 2
    palette: str = "white"
    content_bias: float = 0.5
    # Morphology weights — what fraction of hairs are loops / eyelashes /
    # fragments (the rest are the default curved strand). Defaults give a
    # visible-but-not-dominant mix of variety.
    loop_chance: float = 0.15
    eyelash_chance: float = 0.08
    fragment_chance: float = 0.08
    # Pubic / body-hair morphology — low default so it's a surprise, not a theme.
    kink_chance: float = 0.06
    # Probability that an already-placed hair gets a "buddy" placed nearby,
    # rolled repeatedly per hair (so a single hair can be the head of a small
    # tuft of 1–3). Real shed hair clumps; even spread looks artificial.
    cluster_chance: float = 0.22
    # Always populated after construction (see __post_init__) so callers can
    # echo the seed back to the user for "re-strand with this seed".
    seed: int = field(default_factory=lambda: random.randrange(1, 2**31 - 1))

    def __post_init__(self):
        if self.seed is None:  # accept None defensively
            self.seed = random.randrange(1, 2**31 - 1)

    def rng(self) -> random.Random:
        return random.Random(self.seed)


def _normalize_intensity(intensity: str) -> str:
    raw = (intensity or "").lower().strip()
    if raw not in INTENSITY_TIERS:
        return DEFAULT_INTENSITY
    return raw


def options_from_ui(palette: str, intensity: str, seed: int | None = None) -> InjectOptions:
    """Translate the UI's coarse choices into InjectOptions."""
    palette = (palette or "").lower().strip()
    if palette not in PALETTE_NAMES:
        palette = DEFAULT_PALETTE
    tier = INTENSITY_TIERS[_normalize_intensity(intensity)]
    kwargs = dict(
        rate=tier["rate"],
        image_count=tier["image_count"],
        hairs_per_page=tier["hairs_per_page"],
        palette=palette,
    )
    if seed is not None:
        kwargs["seed"] = seed
    return InjectOptions(**kwargs)


def inject_image_bytes(data: bytes, suffix: str, opts: InjectOptions) -> tuple[bytes, dict]:
    """Overlay hairs onto an image. Returns (encoded bytes, stats dict)."""
    suffix = suffix.lower()
    if suffix not in IMAGE_SUFFIXES:
        raise ValueError(f"unsupported image suffix: {suffix}")

    rng = opts.rng()
    stats = _empty_stats()
    stats["pages_touched"] = 1
    with Image.open(io.BytesIO(data)) as src:
        img = src.convert("RGBA")

    iw, ih = img.size
    regions = _image_content_regions(img)
    units_per_cm = _image_pixels_per_cm(img)
    dpi_info = img.info.get("dpi")
    dpi_value = int(round(min(_MAX_IMAGE_DPI, float(dpi_info[0])))) if (
        dpi_info and isinstance(dpi_info, tuple) and dpi_info[0]
    ) else _DEFAULT_DPI
    stats["substrate"] = {
        "width_native": iw,
        "height_native": ih,
        "native_unit": "px",
        "dpi": dpi_value,
        "width_cm": round(iw / units_per_cm, 1),
        "height_cm": round(ih / units_per_cm, 1),
    }

    for _ in range(max(1, opts.image_count)):
        hair, morph, color = generate_hair_with_morphology(rng, **_morph_kwargs(opts))
        hair = _rotate_hair(hair, rng)
        stats["hairs"] += 1
        stats["morphologies"][morph] += 1
        _bump(stats["palettes"], color)
        hw, hh = hair.size
        x, y, dw, dh, on_content = _place(
            rng, iw, ih, hw, hh,
            MORPHOLOGY_LENGTH_CM[morph], units_per_cm,
            regions=regions, content_bias=opts.content_bias,
        )
        stats["hair_lengths_cm"].append(round(max(dw, dh) / units_per_cm, 2))
        if on_content:
            stats["content_hits"] += 1
        hair_sized = hair.resize((max(1, int(dw)), max(1, int(dh))), Image.LANCZOS)
        img.alpha_composite(hair_sized, (int(x), int(y)))

        buddy_count = 0
        for bx, by in _buddy_offsets(rng, dw, dh, iw, ih, x, y, opts.cluster_chance):
            buddy, bmorph, bcolor = generate_hair_with_morphology(rng, **_morph_kwargs(opts))
            buddy = _rotate_hair(buddy, rng)
            stats["hairs"] += 1
            stats["morphologies"][bmorph] += 1
            _bump(stats["palettes"], bcolor)
            buddy_count += 1
            scale = rng.uniform(0.7, 1.0)
            bw, bh = max(1, int(dw * scale)), max(1, int(dh * scale))
            stats["hair_lengths_cm"].append(round(max(bw, bh) / units_per_cm, 2))
            buddy = buddy.resize((bw, bh), Image.LANCZOS)
            img.alpha_composite(buddy, (int(bx), int(by)))
        if buddy_count > 0:
            stats["clusters"] += 1

    out = io.BytesIO()
    if suffix in (".jpg", ".jpeg"):
        img.convert("RGB").save(out, format="JPEG", quality=92, optimize=True)
    elif suffix == ".gif":
        img.convert("P", palette=Image.ADAPTIVE).save(out, format="GIF")
    elif suffix == ".bmp":
        img.convert("RGB").save(out, format="BMP")
    elif suffix == ".webp":
        img.save(out, format="WEBP", quality=92)
    elif suffix in (".tif", ".tiff"):
        img.save(out, format="TIFF")
    elif suffix == ".ico":
        # ICO frames cap at 256px; let Pillow pick the embedded sizes from the
        # (possibly downscaled) source so the save doesn't fail on large inputs.
        img.save(out, format="ICO")
    else:
        img.save(out, format="PNG")
    return out.getvalue(), stats


def inject_pdf_bytes(data: bytes, opts: InjectOptions) -> tuple[bytes, dict]:
    """Overlay hairs onto a PDF. Returns (rewritten bytes, stats dict)."""
    import fitz  # PyMuPDF

    rng = opts.rng()
    stats = _empty_stats()
    doc = fitz.open(stream=data, filetype="pdf")

    try:
        page_count = doc.page_count
        if page_count == 0:
            return data, stats

        # Report the first page's dimensions as the substrate. Mixed-size PDFs
        # are rare and we'd rather show one honest size than an "it varies".
        first = doc[0]
        stats["substrate"] = {
            "width_native": round(first.rect.width, 1),
            "height_native": round(first.rect.height, 1),
            "native_unit": "pt",
            "dpi": None,
            "width_cm": round(first.rect.width / _POINTS_PER_CM, 1),
            "height_cm": round(first.rect.height / _POINTS_PER_CM, 1),
            "page_count": page_count,
        }

        # Pick which pages get hair. Per-page Bernoulli with at-least-one guarantee
        # so a 2-page doc on subtle still gets visibly haired.
        page_hit = [rng.random() < opts.rate for _ in range(page_count)]
        if not any(page_hit):
            page_hit[rng.randrange(page_count)] = True

        # Record the first page that picked up hair so the UI can render a
        # meaningful before/after preview — showing a blank page from a 50-page
        # PDF would defeat the whole point.
        stats["preview_page"] = next(i for i, hit in enumerate(page_hit) if hit)

        per_page = max(1, opts.hairs_per_page)
        for i, do_it in enumerate(page_hit):
            if not do_it:
                continue
            stats["pages_touched"] += 1
            page = doc[i]
            regions = _pdf_content_regions(page)
            pw = page.rect.width
            ph = page.rect.height
            for _ in range(per_page):
                hair, morph, color = generate_hair_with_morphology(rng, **_morph_kwargs(opts))
                hair = _rotate_hair(hair, rng)
                stats["hairs"] += 1
                stats["morphologies"][morph] += 1
                _bump(stats["palettes"], color)
                hw, hh = hair.size
                x, y, dw, dh, on_content = _place(
                    rng, pw, ph, hw, hh,
                    MORPHOLOGY_LENGTH_CM[morph], _POINTS_PER_CM,
                    regions=regions, content_bias=opts.content_bias,
                )
                stats["hair_lengths_cm"].append(round(max(dw, dh) / _POINTS_PER_CM, 2))
                if on_content:
                    stats["content_hits"] += 1

                buf = io.BytesIO()
                hair.save(buf, format="PNG")
                hair_png = buf.getvalue()
                page.insert_image(
                    fitz.Rect(x, y, x + dw, y + dh),
                    stream=hair_png,
                    keep_proportion=False,
                    overlay=True,
                )

                buddy_count = 0
                for bx, by in _buddy_offsets(rng, dw, dh, pw, ph, x, y, opts.cluster_chance):
                    buddy, bmorph, bcolor = generate_hair_with_morphology(rng, **_morph_kwargs(opts))
                    buddy = _rotate_hair(buddy, rng)
                    stats["hairs"] += 1
                    stats["morphologies"][bmorph] += 1
                    _bump(stats["palettes"], bcolor)
                    buddy_count += 1
                    bs = rng.uniform(0.7, 1.0)
                    bw_, bh_ = dw * bs, dh * bs
                    stats["hair_lengths_cm"].append(round(max(bw_, bh_) / _POINTS_PER_CM, 2))
                    bbuf = io.BytesIO()
                    buddy.save(bbuf, format="PNG")
                    page.insert_image(
                        fitz.Rect(bx, by, bx + bw_, by + bh_),
                        stream=bbuf.getvalue(),
                        keep_proportion=False,
                        overlay=True,
                    )
                if buddy_count > 0:
                    stats["clusters"] += 1

        return doc.tobytes(garbage=4, deflate=True), stats
    finally:
        doc.close()


def render_pdf_page_png(pdf_bytes: bytes, page_index: int = 0, max_pixels: int = 1400) -> bytes:
    """Rasterise one PDF page to PNG. Used to build the before/after preview
    pair shown in the in-browser result panel — the page that picked up the
    first round of hairs is the one worth showing. PyMuPDF renders a page in
    a few ms even at preview resolution. `max_pixels` caps the larger output
    dimension so the preview reads cleanly in the narrow result column."""
    import fitz
    doc = fitz.open(stream=pdf_bytes, filetype="pdf")
    try:
        if page_index < 0 or page_index >= doc.page_count:
            page_index = 0
        page = doc[page_index]
        w = max(page.rect.width, 1.0)
        h = max(page.rect.height, 1.0)
        scale = min(max_pixels / w, max_pixels / h)
        pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale), alpha=False)
        return pix.tobytes("png")
    finally:
        doc.close()


def inject_pptx_bytes(data: bytes, opts: InjectOptions) -> tuple[bytes, dict]:
    """Overlay hairs onto a pptx. Returns (rewritten bytes, stats dict)."""
    from pptx import Presentation

    rng = opts.rng()
    stats = _empty_stats()
    prs = Presentation(io.BytesIO(data))
    sw, sh = prs.slide_width, prs.slide_height  # EMUs

    slides = list(prs.slides)
    if not slides:
        out = io.BytesIO()
        prs.save(out)
        return out.getvalue(), stats

    stats["substrate"] = {
        "width_native": sw,
        "height_native": sh,
        "native_unit": "EMU",
        "dpi": None,
        "width_cm": round(sw / _EMU_PER_CM, 1),
        "height_cm": round(sh / _EMU_PER_CM, 1),
        "slide_count": len(slides),
    }

    slide_hit = [rng.random() < opts.rate for _ in slides]
    if not any(slide_hit):
        slide_hit[rng.randrange(len(slides))] = True

    per_slide = max(1, opts.hairs_per_page)
    for slide, do_it in zip(slides, slide_hit):
        if not do_it:
            continue
        stats["pages_touched"] += 1
        regions = _pptx_content_regions(slide)
        for _ in range(per_slide):
            hair, morph, color = generate_hair_with_morphology(rng, **_morph_kwargs(opts))
            hair = _rotate_hair(hair, rng)
            stats["hairs"] += 1
            stats["morphologies"][morph] += 1
            _bump(stats["palettes"], color)
            hw, hh = hair.size
            x, y, dw, dh, on_content = _place(
                rng, sw, sh, hw, hh,
                MORPHOLOGY_LENGTH_CM[morph], _EMU_PER_CM,
                regions=regions, content_bias=opts.content_bias,
            )
            stats["hair_lengths_cm"].append(round(max(dw, dh) / _EMU_PER_CM, 2))
            if on_content:
                stats["content_hits"] += 1

            buf = io.BytesIO()
            hair.save(buf, format="PNG")
            buf.seek(0)
            slide.shapes.add_picture(buf, int(x), int(y), width=int(dw), height=int(dh))

            buddy_count = 0
            for bx, by in _buddy_offsets(rng, dw, dh, sw, sh, x, y, opts.cluster_chance):
                buddy, bmorph, bcolor = generate_hair_with_morphology(rng, **_morph_kwargs(opts))
                buddy = _rotate_hair(buddy, rng)
                stats["hairs"] += 1
                stats["morphologies"][bmorph] += 1
                _bump(stats["palettes"], bcolor)
                buddy_count += 1
                bs = rng.uniform(0.7, 1.0)
                stats["hair_lengths_cm"].append(
                    round(max(dw, dh) * bs / _EMU_PER_CM, 2)
                )
                bbuf = io.BytesIO()
                buddy.save(bbuf, format="PNG")
                bbuf.seek(0)
                slide.shapes.add_picture(
                    bbuf, int(bx), int(by),
                    width=int(dw * bs), height=int(dh * bs),
                )
            if buddy_count > 0:
                stats["clusters"] += 1

    out = io.BytesIO()
    prs.save(out)
    return out.getvalue(), stats


# wp:anchor template for a page-relative floating picture in a .docx. The
# graphic subtree (which python-docx builds with the correct image relationship)
# is moved in as the final child; everything else is positioning. behindDoc=1 +
# wrapNone lets the hair sit over the text the way a stray hair sits over print.
_DOCX_ANCHOR_XML = (
    '<wp:anchor xmlns:wp="http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing"'
    ' distT="0" distB="0" distL="0" distR="0" simplePos="0" relativeHeight="2"'
    ' behindDoc="1" locked="0" layoutInCell="1" allowOverlap="1">'
    '<wp:simplePos x="0" y="0"/>'
    '<wp:positionH relativeFrom="page"><wp:posOffset>{x}</wp:posOffset></wp:positionH>'
    '<wp:positionV relativeFrom="page"><wp:posOffset>{y}</wp:posOffset></wp:positionV>'
    '<wp:extent cx="{cx}" cy="{cy}"/>'
    '<wp:effectExtent l="0" t="0" r="0" b="0"/>'
    '<wp:wrapNone/>'
    '<wp:docPr id="{id}" name="StrandHair{id}"/>'
    '<wp:cNvGraphicFramePr/>'
    '</wp:anchor>'
)


def _docx_float_picture(paragraph, png_bytes: bytes, x_emu: int, y_emu: int,
                        cx_emu: int, cy_emu: int, drawing_id: int) -> None:
    """Add one page-anchored floating picture to `paragraph`.

    python-docx only exposes inline pictures, so we let it build the inline
    drawing (which wires up the image part + relationship correctly), then
    swap the `wp:inline` wrapper for a `wp:anchor` carrying the same graphic.
    """
    from docx.oxml import parse_xml
    from docx.oxml.ns import qn
    from docx.shared import Emu

    run = paragraph.add_run()
    run.add_picture(io.BytesIO(png_bytes), width=Emu(int(cx_emu)), height=Emu(int(cy_emu)))

    drawing = run._r.find(qn("w:drawing"))
    inline = drawing[0]
    graphic = inline.find(qn("a:graphic"))

    anchor = parse_xml(_DOCX_ANCHOR_XML.format(
        x=int(x_emu), y=int(y_emu), cx=int(cx_emu), cy=int(cy_emu), id=drawing_id,
    ))
    anchor.append(graphic)  # lxml moves the node; correct schema order (last child)
    drawing.replace(inline, anchor)


# Default Word body metrics for the page estimate below. A floating image
# anchored to a paragraph renders on whatever page that paragraph flows onto,
# so to spread hairs across the document the way the PDF/pptx injectors spread
# across pages/slides we first estimate which page each paragraph lands on.
_EMU_PER_PT = 12700
_DOCX_LINE_EMU = int(11 * 1.15 * _EMU_PER_PT)   # 11pt body text, 1.15 line spacing
_DOCX_CHAR_EMU = int(5.0 * _EMU_PER_PT)         # ~half the point size per glyph


def _docx_page_breaks(para) -> int:
    """How many page boundaries this paragraph forces.

    Counts hard breaks (`<w:br w:type="page"/>`) and Word's own soft
    `lastRenderedPageBreak` hints — the latter is authoritative for documents
    last saved by Word, the former covers ones built programmatically.
    """
    xml = para._p.xml
    return xml.count('w:type="page"') + xml.count("lastRenderedPageBreak")


def _docx_estimate_pages(paras, content_w_emu: int, content_h_emu: int):
    """Estimate a page index for each paragraph; return (page_of, n_pages).

    A coarse text-flow model: each paragraph occupies ceil(chars / chars-per-
    line) lines of body height, plus any explicit/soft page breaks. Good enough
    to scatter hairs across the document — it doesn't need to match Word's
    layout engine exactly, just to put different hairs on different pages.
    """
    cpl = max(1, content_w_emu // _DOCX_CHAR_EMU)
    page = 0
    cum = 0
    page_of = []
    for p in paras:
        breaks = _docx_page_breaks(p)
        if breaks:
            page += breaks
            cum = 0
        page_of.append(page)
        text_len = len(p.text or "")
        lines = max(1, -(-text_len // cpl))  # ceil division
        cum += lines * _DOCX_LINE_EMU
        while cum >= content_h_emu:
            page += 1
            cum -= content_h_emu
    return page_of, page + 1


def inject_docx_bytes(data: bytes, opts: InjectOptions) -> tuple[bytes, dict]:
    """Overlay hairs onto a Word document. Returns (rewritten bytes, stats dict).

    Word reflows text across pages at render time, so there's no page grid to
    walk the way a PDF has. We estimate which page each paragraph lands on, then
    use the same per-page `rate` + `hairs_per_page` model as the PDF/pptx
    injectors — each hair is anchored (page-relative) to a paragraph on its
    target page, so Word renders the hairs spread through the document.
    """
    from docx import Document

    rng = opts.rng()
    stats = _empty_stats()
    doc = Document(io.BytesIO(data))

    section = doc.sections[0]
    pw = int(section.page_width)   # EMU
    ph = int(section.page_height)  # EMU
    content_w = max(1, pw - int(section.left_margin or 0) - int(section.right_margin or 0))
    content_h = max(1, ph - int(section.top_margin or 0) - int(section.bottom_margin or 0))

    paras = doc.paragraphs or [doc.add_paragraph()]
    page_of, n_pages = _docx_estimate_pages(paras, content_w, content_h)

    stats["substrate"] = {
        "width_native": pw,
        "height_native": ph,
        "native_unit": "EMU",
        "dpi": None,
        "width_cm": round(pw / _EMU_PER_CM, 1),
        "height_cm": round(ph / _EMU_PER_CM, 1),
        "page_count": n_pages,  # estimated — Word's exact pagination may differ
    }

    # Group paragraphs by estimated page. A page-break paragraph is excluded as
    # an anchor candidate: its anchor character straddles the boundary, so a
    # float pinned to it renders on the *previous* page in practice — anchoring
    # to a real text paragraph on the page is what reliably lands the hair there.
    paras_by_page: dict[int, list] = {}
    all_by_page: dict[int, list] = {}
    for para, pg in zip(paras, page_of):
        all_by_page.setdefault(pg, []).append(para)
        if _docx_page_breaks(para) or not (para.text or "").strip():
            continue
        paras_by_page.setdefault(pg, []).append(para)
    # If filtering left a page with no text anchor, fall back to any paragraph
    # the estimate placed there (e.g. a page of only images or breaks).
    for pg, plist in all_by_page.items():
        paras_by_page.setdefault(pg, plist)

    # Per-page Bernoulli with an at-least-one guarantee, exactly like the PDF
    # path — a short doc on "subtle" still picks up a visible hair.
    page_hit = [rng.random() < opts.rate for _ in range(n_pages)]
    if not any(page_hit):
        page_hit[rng.randrange(n_pages)] = True

    def anchors_for(pg: int):
        # Fall back to the nearest populated page if the estimate left this one
        # empty (e.g. a forced break with no following paragraph yet).
        if pg in paras_by_page:
            return paras_by_page[pg]
        nearest = min(paras_by_page, key=lambda k: abs(k - pg))
        return paras_by_page[nearest]

    drawing_id = 1
    per_page = max(1, opts.hairs_per_page)
    for pg, hit in enumerate(page_hit):
        if not hit:
            continue
        stats["pages_touched"] += 1
        candidates = anchors_for(pg)
        for _ in range(per_page):
            anchor_para = rng.choice(candidates)
            hair, morph, color = generate_hair_with_morphology(rng, **_morph_kwargs(opts))
            hair = _rotate_hair(hair, rng)
            stats["hairs"] += 1
            stats["morphologies"][morph] += 1
            _bump(stats["palettes"], color)
            hw, hh = hair.size
            x, y, dw, dh, on_content = _place(
                rng, pw, ph, hw, hh,
                MORPHOLOGY_LENGTH_CM[morph], _EMU_PER_CM,
            )
            stats["hair_lengths_cm"].append(round(max(dw, dh) / _EMU_PER_CM, 2))
            if on_content:
                stats["content_hits"] += 1

            buf = io.BytesIO()
            hair.save(buf, format="PNG")
            _docx_float_picture(anchor_para, buf.getvalue(), x, y, dw, dh, drawing_id)
            drawing_id += 1

            buddy_count = 0
            for bx, by in _buddy_offsets(rng, dw, dh, pw, ph, x, y, opts.cluster_chance):
                buddy, bmorph, bcolor = generate_hair_with_morphology(rng, **_morph_kwargs(opts))
                buddy = _rotate_hair(buddy, rng)
                stats["hairs"] += 1
                stats["morphologies"][bmorph] += 1
                _bump(stats["palettes"], bcolor)
                buddy_count += 1
                bs = rng.uniform(0.7, 1.0)
                bw_, bh_ = dw * bs, dh * bs
                stats["hair_lengths_cm"].append(round(max(bw_, bh_) / _EMU_PER_CM, 2))
                bbuf = io.BytesIO()
                buddy.save(bbuf, format="PNG")
                _docx_float_picture(anchor_para, bbuf.getvalue(), bx, by, bw_, bh_, drawing_id)
                drawing_id += 1
            if buddy_count > 0:
                stats["clusters"] += 1

    out = io.BytesIO()
    doc.save(out)
    return out.getvalue(), stats


# A spreadsheet has no page; we synthesise a canvas from the used range using
# Excel's default cell sizing so hairs scatter over the populated area rather
# than a single cell. Approximate by design — it's a novelty overlay.
_XLSX_EMU_PER_COL = 609600   # default column width ≈ 64 px
_XLSX_EMU_PER_ROW = 190500   # default row height = 15 pt ≈ 20 px


def inject_xlsx_bytes(data: bytes, opts: InjectOptions) -> tuple[bytes, dict]:
    """Overlay hairs onto an Excel workbook. Returns (rewritten bytes, stats dict).

    Each selected sheet gets hairs absolutely positioned (EMU) over a canvas
    sized from its used range. Per-sheet selection mirrors the PDF page model.
    """
    from openpyxl import load_workbook
    from openpyxl.drawing.image import Image as XLImage
    from openpyxl.drawing.spreadsheet_drawing import AbsoluteAnchor
    from openpyxl.drawing.xdr import XDRPoint2D, XDRPositiveSize2D

    rng = opts.rng()
    stats = _empty_stats()
    wb = load_workbook(io.BytesIO(data))

    sheets = wb.worksheets
    if not sheets:
        out = io.BytesIO()
        wb.save(out)
        return out.getvalue(), stats

    def canvas_emu(ws):
        cols = max(ws.max_column or 1, 8)
        rows = max(ws.max_row or 1, 20)
        return cols * _XLSX_EMU_PER_COL, rows * _XLSX_EMU_PER_ROW

    cw0, ch0 = canvas_emu(sheets[0])
    stats["substrate"] = {
        "width_native": cw0,
        "height_native": ch0,
        "native_unit": "EMU",
        "dpi": None,
        "width_cm": round(cw0 / _EMU_PER_CM, 1),
        "height_cm": round(ch0 / _EMU_PER_CM, 1),
        "sheet_count": len(sheets),
    }

    sheet_hit = [rng.random() < opts.rate for _ in sheets]
    if not any(sheet_hit):
        sheet_hit[rng.randrange(len(sheets))] = True

    def add_hair(ws, png_bytes, x, y, cx, cy):
        img = XLImage(io.BytesIO(png_bytes))
        img.anchor = AbsoluteAnchor(
            pos=XDRPoint2D(int(x), int(y)),
            ext=XDRPositiveSize2D(int(cx), int(cy)),
        )
        ws.add_image(img)

    per_sheet = max(1, opts.hairs_per_page)
    for ws, do_it in zip(sheets, sheet_hit):
        if not do_it:
            continue
        stats["pages_touched"] += 1
        cw, ch = canvas_emu(ws)
        for _ in range(per_sheet):
            hair, morph, color = generate_hair_with_morphology(rng, **_morph_kwargs(opts))
            hair = _rotate_hair(hair, rng)
            stats["hairs"] += 1
            stats["morphologies"][morph] += 1
            _bump(stats["palettes"], color)
            hw, hh = hair.size
            x, y, dw, dh, on_content = _place(
                rng, cw, ch, hw, hh,
                MORPHOLOGY_LENGTH_CM[morph], _EMU_PER_CM,
            )
            stats["hair_lengths_cm"].append(round(max(dw, dh) / _EMU_PER_CM, 2))
            if on_content:
                stats["content_hits"] += 1

            buf = io.BytesIO()
            hair.save(buf, format="PNG")
            add_hair(ws, buf.getvalue(), x, y, dw, dh)

            buddy_count = 0
            for bx, by in _buddy_offsets(rng, dw, dh, cw, ch, x, y, opts.cluster_chance):
                buddy, bmorph, bcolor = generate_hair_with_morphology(rng, **_morph_kwargs(opts))
                buddy = _rotate_hair(buddy, rng)
                stats["hairs"] += 1
                stats["morphologies"][bmorph] += 1
                _bump(stats["palettes"], bcolor)
                buddy_count += 1
                bs = rng.uniform(0.7, 1.0)
                bw_, bh_ = dw * bs, dh * bs
                stats["hair_lengths_cm"].append(round(max(bw_, bh_) / _EMU_PER_CM, 2))
                bbuf = io.BytesIO()
                buddy.save(bbuf, format="PNG")
                add_hair(ws, bbuf.getvalue(), bx, by, bw_, bh_)
            if buddy_count > 0:
                stats["clusters"] += 1

    out = io.BytesIO()
    wb.save(out)
    return out.getvalue(), stats


def _suffix_of(name: str) -> str:
    return ("." + name.rsplit(".", 1)[-1].lower()) if "." in name else ""


def strand_bytes(data: bytes, filename: str, opts: InjectOptions) -> tuple[bytes, dict]:
    """Dispatch to the right injector based on filename suffix.

    Returns (rewritten bytes, stats dict). Stats include total hair count,
    a per-morphology breakdown, and pages touched.
    """
    suffix = _suffix_of(filename)
    if suffix in IMAGE_SUFFIXES:
        return inject_image_bytes(data, suffix, opts)
    if suffix in PDF_SUFFIXES:
        return inject_pdf_bytes(data, opts)
    if suffix in PPTX_SUFFIXES:
        return inject_pptx_bytes(data, opts)
    if suffix in DOCX_SUFFIXES:
        return inject_docx_bytes(data, opts)
    if suffix in XLSX_SUFFIXES:
        return inject_xlsx_bytes(data, opts)
    raise ValueError(f"unsupported file type: {suffix or filename!r}")


def _apply_suffix(name: str, suffix: str) -> str:
    """Insert `suffix` between stem and extension. Empty suffix is a no-op."""
    if not suffix:
        return name
    if "." in name:
        stem, _, ext = name.rpartition(".")
        return f"{stem}{suffix}.{ext}"
    return f"{name}{suffix}"


def strand_zip_bytes(
    data: bytes,
    opts: InjectOptions,
    name_suffix: str = "-strand",
) -> tuple[bytes, dict]:
    """
    Walk every entry in a zip, strand supported files, return (zip_bytes, report).

    Each entry's status is one of: "hairified", "skipped", "errored". A
    `_strand-report.txt` is added to the output zip alongside the entries.
    Directory entries are preserved as-is.
    """
    import zipfile

    base_seed = opts.seed
    report = {
        "entries": [],
        "haired_count": 0,
        "skipped_count": 0,
        "error_count": 0,
        "seed": base_seed,
        "stats": _empty_stats(),
    }

    out_buf = io.BytesIO()
    with zipfile.ZipFile(io.BytesIO(data)) as zin, \
         zipfile.ZipFile(out_buf, "w", compression=zipfile.ZIP_DEFLATED) as zout:

        infos = zin.infolist()
        for idx, info in enumerate(infos):
            name = info.filename
            # Skip our own report from a previous round-trip, and macOS metadata.
            if name == "_strand-report.txt" or name.startswith("__MACOSX/"):
                continue
            if info.is_dir():
                zout.writestr(info, b"")
                continue

            suffix = _suffix_of(name)

            if suffix not in SUPPORTED_SUFFIXES:
                # Pass through unchanged; record as skipped.
                try:
                    raw = zin.read(info)
                    zout.writestr(info, raw)
                except Exception as exc:
                    report["error_count"] += 1
                    report["entries"].append({
                        "name": name, "out_name": None,
                        "status": "errored", "reason": f"read failed: {exc.__class__.__name__}",
                    })
                    continue
                report["skipped_count"] += 1
                report["entries"].append({
                    "name": name, "out_name": name,
                    "status": "skipped",
                    "reason": f"unsupported type ({suffix or 'no extension'})",
                })
                continue

            # Give each entry its own derived seed so re-running with the same
            # outer seed reproduces every inner result, while inner entries
            # still differ from one another.
            per_entry_seed = (base_seed + idx * 1_000_003) & 0x7FFFFFFF
            entry_opts = InjectOptions(
                rate=opts.rate,
                hairs_per_page=opts.hairs_per_page,
                image_count=opts.image_count,
                palette=opts.palette,
                content_bias=opts.content_bias,
                loop_chance=opts.loop_chance,
                eyelash_chance=opts.eyelash_chance,
                fragment_chance=opts.fragment_chance,
                kink_chance=opts.kink_chance,
                cluster_chance=opts.cluster_chance,
                seed=per_entry_seed,
            )

            try:
                raw = zin.read(info)
                processed, entry_stats = strand_bytes(raw, name, entry_opts)
            except Exception as exc:
                report["error_count"] += 1
                report["entries"].append({
                    "name": name, "out_name": None,
                    "status": "errored", "reason": f"{exc.__class__.__name__}: {exc}",
                })
                continue

            _merge_stats(report["stats"], entry_stats)
            out_name = _apply_suffix(name, name_suffix)
            zout.writestr(out_name, processed)
            report["haired_count"] += 1
            report["entries"].append({
                "name": name, "out_name": out_name,
                "status": "hairified", "reason": None,
            })

        # Write the report as the last entry.
        zout.writestr("_strand-report.txt", _format_zip_report(report, opts))

    return out_buf.getvalue(), report


def _format_zip_report(report: dict, opts: InjectOptions) -> str:
    from datetime import datetime, timezone
    lines = [
        "Strand report",
        f"generated: {datetime.now(timezone.utc).isoformat(timespec='seconds')}",
        f"seed:      {report['seed']}",
        f"palette:   {opts.palette}",
        f"rate:      {opts.rate}",
        f"per page:  {opts.hairs_per_page}",
        f"per image: {opts.image_count}",
        "",
        f"hairified: {report['haired_count']}",
        f"skipped:   {report['skipped_count']}",
        f"errored:   {report['error_count']}",
        "",
        "Entries:",
    ]
    for e in report["entries"]:
        if e["status"] == "hairified":
            lines.append(f"  [hairified] {e['name']} -> {e['out_name']}")
        elif e["status"] == "skipped":
            lines.append(f"  [skipped]   {e['name']}  ({e['reason']})")
        else:
            lines.append(f"  [error]     {e['name']}  ({e['reason']})")
    return "\n".join(lines) + "\n"
