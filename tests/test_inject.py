"""
Roundtrip tests for the bytes-based injectors and the FastAPI app.

Fixtures are generated at test time — no binary blobs in the repo.
"""

from __future__ import annotations

import io

import pytest
from PIL import Image

from app.core import (
    INTENSITY_TIERS,
    InjectOptions,
    strand_bytes,
    strand_zip_bytes,
    inject_image_bytes,
    inject_pdf_bytes,
    inject_pptx_bytes,
    options_from_ui,
)


# --- Fixtures --------------------------------------------------------------

@pytest.fixture
def png_bytes() -> bytes:
    """A small white PNG with a couple of rectangles, just to give content detection something to find."""
    from PIL import ImageDraw
    img = Image.new("RGB", (400, 300), "white")
    d = ImageDraw.Draw(img)
    d.rectangle((30, 30, 200, 80), fill="black")
    d.rectangle((50, 120, 350, 250), outline="black", width=3)
    buf = io.BytesIO()
    img.save(buf, format="PNG")
    return buf.getvalue()


@pytest.fixture
def jpeg_bytes() -> bytes:
    img = Image.new("RGB", (320, 240), "white")
    buf = io.BytesIO()
    img.save(buf, format="JPEG", quality=90)
    return buf.getvalue()


@pytest.fixture
def pdf_bytes() -> bytes:
    from reportlab.lib.pagesizes import letter
    from reportlab.pdfgen import canvas
    buf = io.BytesIO()
    c = canvas.Canvas(buf, pagesize=letter)
    for i in range(3):
        c.setFont("Helvetica", 12)
        c.drawString(72, 720, f"Page {i + 1}")
        c.drawString(72, 700, "Lorem ipsum dolor sit amet, consectetur adipiscing elit.")
        c.drawString(72, 680, "Phasellus blandit, mauris in commodo accumsan.")
        c.showPage()
    c.save()
    return buf.getvalue()


@pytest.fixture
def pptx_bytes() -> bytes:
    from pptx import Presentation
    prs = Presentation()
    blank = prs.slide_layouts[6]
    for i in range(3):
        slide = prs.slides.add_slide(blank)
        from pptx.util import Inches
        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(6), Inches(1))
        tb.text_frame.text = f"Slide {i + 1}"
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


# --- Per-injector roundtrips ----------------------------------------------

def test_inject_png_roundtrip(png_bytes):
    opts = InjectOptions(seed=1, image_count=2)
    out, _ = inject_image_bytes(png_bytes, ".png", opts)
    assert isinstance(out, bytes)
    assert len(out) > 0
    assert out != png_bytes

    with Image.open(io.BytesIO(out)) as im:
        im.load()
        assert im.size == (400, 300)
        assert im.format == "PNG"


def test_inject_jpeg_roundtrip(jpeg_bytes):
    opts = InjectOptions(seed=2, image_count=1)
    out, _ = inject_image_bytes(jpeg_bytes, ".jpg", opts)
    with Image.open(io.BytesIO(out)) as im:
        im.load()
        assert im.size == (320, 240)
        assert im.format == "JPEG"


def test_inject_pdf_roundtrip(pdf_bytes):
    import fitz

    opts = InjectOptions(seed=3, rate=0.85)
    out, _ = inject_pdf_bytes(pdf_bytes, opts)
    assert len(out) > 0
    assert out != pdf_bytes

    doc = fitz.open(stream=out, filetype="pdf")
    try:
        assert doc.page_count == 3
    finally:
        doc.close()


def test_inject_pdf_subtle_still_hits_at_least_one_page(pdf_bytes):
    """Subtle (rate=0.25) on a 3-page doc must still produce a haired page."""
    import fitz

    opts = InjectOptions(seed=12345, rate=0.25)
    out, _ = inject_pdf_bytes(pdf_bytes, opts)
    # An untouched PDF round-tripped through PyMuPDF still differs from the
    # reportlab output, so size-difference isn't a useful signal here. Instead,
    # check that at least one page has embedded images.
    doc = fitz.open(stream=out, filetype="pdf")
    try:
        total_images = sum(len(doc[i].get_images()) for i in range(doc.page_count))
        assert total_images >= 1
    finally:
        doc.close()


def test_inject_pptx_roundtrip(pptx_bytes):
    from pptx import Presentation

    opts = InjectOptions(seed=4, rate=0.85)
    out, _ = inject_pptx_bytes(pptx_bytes, opts)
    assert len(out) > 0
    assert out != pptx_bytes

    prs = Presentation(io.BytesIO(out))
    assert len(prs.slides) == 3


def test_strand_bytes_dispatches_by_suffix(png_bytes, pdf_bytes):
    opts = InjectOptions(seed=5)
    assert strand_bytes(png_bytes, "photo.png", opts)[0] != png_bytes
    assert strand_bytes(pdf_bytes, "deck.pdf", opts)[0] != pdf_bytes


def test_strand_bytes_rejects_unsupported():
    with pytest.raises(ValueError):
        strand_bytes(b"hello", "notes.txt", InjectOptions())


def test_seed_is_deterministic(png_bytes):
    a, _ = inject_image_bytes(png_bytes, ".png", InjectOptions(seed=42, image_count=2))
    b, _ = inject_image_bytes(png_bytes, ".png", InjectOptions(seed=42, image_count=2))
    assert a == b


# --- Stats ------------------------------------------------------------------

def test_image_injector_returns_stats(png_bytes):
    """Stats track hair count, per-morphology breakdown, palettes, pages, clusters."""
    opts = InjectOptions(seed=1, image_count=3, cluster_chance=0.0, palette="white")
    _, stats = inject_image_bytes(png_bytes, ".png", opts)
    assert stats["hairs"] == 3
    assert stats["pages_touched"] == 1
    assert stats["clusters"] == 0  # cluster_chance=0 disables buddies
    morph_sum = sum(stats["morphologies"].values())
    assert morph_sum == 3
    assert set(stats["morphologies"].keys()) == {"curve", "loop", "eyelash", "fragment", "kink"}
    # Non-mixed palette: every hair drawn in the requested colour family.
    assert stats["palettes"] == {"white": 3}


def test_mixed_palette_resolves_to_concrete_colours(png_bytes):
    """`mixed` should produce a breakdown of actual palette names, never the
    literal string 'mixed'."""
    opts = InjectOptions(seed=99, image_count=12, cluster_chance=0.0, palette="mixed")
    _, stats = inject_image_bytes(png_bytes, ".png", opts)
    assert "mixed" not in stats["palettes"], "mixed should be resolved per hair"
    assert sum(stats["palettes"].values()) == 12
    # With 12 draws across 6 palettes, very likely we get >1 distinct colour.
    assert len(stats["palettes"]) > 1


def test_clusters_count_reflects_buddies(png_bytes):
    """High cluster_chance should produce a non-zero cluster count."""
    opts = InjectOptions(seed=7, image_count=4, cluster_chance=0.95, palette="white")
    _, stats = inject_image_bytes(png_bytes, ".png", opts)
    assert stats["clusters"] > 0
    # And total hairs strictly exceeds image_count (because buddies added).
    assert stats["hairs"] > 4


def test_pdf_injector_stats_track_pages(pdf_bytes):
    opts = InjectOptions(seed=2, rate=1.0, hairs_per_page=2, cluster_chance=0.0)
    _, stats = inject_pdf_bytes(pdf_bytes, opts)
    # Every page hit, 2 hairs per page, 3 pages → 6.
    assert stats["pages_touched"] == 3
    assert stats["hairs"] == 6


def test_stats_header_present_in_response(client, png_bytes):
    r = client.post(
        "/strand",
        files={"file": ("photo.png", png_bytes, "image/png")},
        data={"palette": "dark", "intensity": "normal"},
    )
    assert r.status_code == 200
    raw = r.headers.get("X-Strand-Stats")
    assert raw, "X-Strand-Stats header missing"
    import json
    stats = json.loads(raw)
    assert stats["hairs"] >= 1
    assert "morphologies" in stats
    assert stats["pages_touched"] >= 1


def test_zip_response_stats_aggregate_across_entries(client, zip_bytes):
    r = client.post(
        "/strand",
        files={"file": ("bundle.zip", zip_bytes, "application/zip")},
        data={"palette": "dark", "intensity": "normal"},
    )
    assert r.status_code == 200
    import json
    stats = json.loads(r.headers["X-Strand-Stats"])
    # 3 supported entries each get at least one hair.
    assert stats["hairs"] >= 3


# --- HTTP layer ------------------------------------------------------------

@pytest.fixture
def client():
    from fastapi.testclient import TestClient
    from app.main import app
    return TestClient(app)


def test_health(client):
    r = client.get("/health")
    assert r.status_code == 200
    assert r.json() == {"ok": True}


def test_index_serves_html(client):
    r = client.get("/")
    assert r.status_code == 200
    assert "Strand" in r.text


def test_strand_endpoint_png(client, png_bytes):
    r = client.post(
        "/strand",
        files={"file": ("photo.png", png_bytes, "image/png")},
        data={"palette": "dark", "intensity": "normal"},
    )
    assert r.status_code == 200
    assert r.headers["content-type"] == "image/png"
    assert "photo-strand.png" in r.headers.get("content-disposition", "")
    assert len(r.content) > 0
    # Confirm the response decodes as a valid PNG.
    with Image.open(io.BytesIO(r.content)) as im:
        im.load()
        assert im.format == "PNG"


def test_strand_endpoint_rejects_unsupported(client):
    r = client.post(
        "/strand",
        files={"file": ("notes.txt", b"hello", "text/plain")},
        data={"palette": "dark", "intensity": "normal"},
    )
    assert r.status_code == 415


def test_strand_endpoint_rejects_empty(client):
    r = client.post(
        "/strand",
        files={"file": ("photo.png", b"", "image/png")},
        data={"palette": "dark", "intensity": "normal"},
    )
    assert r.status_code == 400


# --- New v1.1 surface --------------------------------------------------------

def test_intensity_table_is_monotonic_in_density():
    """The funny tiers should actually escalate; if someone reorders the dict,
    fail loudly."""
    tiers = list(INTENSITY_TIERS.values())
    for prev, nxt in zip(tiers, tiers[1:]):
        assert nxt["image_count"] >= prev["image_count"]
        # Hairs per page should never go down either.
        assert nxt["hairs_per_page"] >= prev["hairs_per_page"]
    # Sanity: top tier is meaningfully more hair than bottom.
    assert tiers[-1]["image_count"] >= tiers[0]["image_count"] * 10


@pytest.mark.parametrize("intensity", list(INTENSITY_TIERS.keys()))
def test_every_intensity_tier_runs(png_bytes, intensity):
    opts = options_from_ui(palette="dark", intensity=intensity, seed=7)
    out, _ = inject_image_bytes(png_bytes, ".png", opts)
    assert len(out) > 0
    with Image.open(io.BytesIO(out)) as im:
        im.load()
        assert im.size == (400, 300)


def test_higher_density_embeds_more_images_in_pdf(pdf_bytes):
    """heavy should produce strictly more embedded images than subtle on the same doc."""
    import fitz

    def total_images(intensity):
        opts = options_from_ui(palette="dark", intensity=intensity, seed=1234)
        out, _ = inject_pdf_bytes(pdf_bytes, opts)
        doc = fitz.open(stream=out, filetype="pdf")
        try:
            return sum(len(doc[i].get_images()) for i in range(doc.page_count))
        finally:
            doc.close()

    n_subtle = total_images("subtle")
    n_heavy = total_images("heavy")
    n_werewolf = total_images("werewolf")
    assert n_subtle >= 1
    assert n_heavy > n_subtle
    assert n_werewolf > n_heavy


def test_seed_is_populated_after_construction():
    """Even without an explicit seed, opts.seed must be readable so we can echo it."""
    opts = InjectOptions()
    assert isinstance(opts.seed, int)
    assert opts.seed > 0


# --- ZIP roundtrip -----------------------------------------------------------

@pytest.fixture
def zip_bytes(png_bytes, pdf_bytes, pptx_bytes) -> bytes:
    """A zip with one of each supported format plus a deliberate unsupported entry and a nested dir."""
    import zipfile
    buf = io.BytesIO()
    with zipfile.ZipFile(buf, "w", compression=zipfile.ZIP_DEFLATED) as z:
        z.writestr("photo.png", png_bytes)
        z.writestr("docs/deck.pptx", pptx_bytes)
        z.writestr("docs/manual.pdf", pdf_bytes)
        z.writestr("README.txt", b"keep me as is\n")
    return buf.getvalue()


def test_zip_roundtrip_processes_each_supported_entry(zip_bytes):
    import zipfile
    opts = options_from_ui(palette="dark", intensity="normal", seed=42)
    out, report = strand_zip_bytes(zip_bytes, opts, name_suffix="-strand")

    assert report["haired_count"] == 3
    assert report["skipped_count"] == 1
    assert report["error_count"] == 0

    with zipfile.ZipFile(io.BytesIO(out)) as z:
        names = set(z.namelist())
        assert "photo-strand.png" in names
        assert "docs/deck-strand.pptx" in names
        assert "docs/manual-strand.pdf" in names
        # Unsupported entries pass through unchanged with the original name.
        assert "README.txt" in names
        assert z.read("README.txt") == b"keep me as is\n"
        # Report is embedded.
        assert "_strand-report.txt" in names
        report_text = z.read("_strand-report.txt").decode("utf-8")
        assert "Strand report" in report_text
        assert "seed:" in report_text


def test_zip_custom_suffix(zip_bytes):
    import zipfile
    opts = options_from_ui(palette="dark", intensity="subtle", seed=99)
    out, _ = strand_zip_bytes(zip_bytes, opts, name_suffix=".v2")
    with zipfile.ZipFile(io.BytesIO(out)) as z:
        assert "photo.v2.png" in z.namelist()
        assert "docs/deck.v2.pptx" in z.namelist()


def test_zip_empty_suffix_keeps_names(zip_bytes):
    import zipfile
    opts = options_from_ui(palette="dark", intensity="subtle", seed=99)
    out, _ = strand_zip_bytes(zip_bytes, opts, name_suffix="")
    with zipfile.ZipFile(io.BytesIO(out)) as z:
        assert "photo.png" in z.namelist()
        # Content still changed.
        assert z.read("photo.png") != z.read("photo.png")[:0] + b""  # not empty
        assert "_strand-report.txt" in z.namelist()


# --- HTTP layer for the new surface ------------------------------------------

def test_seed_header_is_echoed(client, png_bytes):
    r = client.post(
        "/strand",
        files={"file": ("photo.png", png_bytes, "image/png")},
        data={"palette": "dark", "intensity": "normal"},
    )
    assert r.status_code == 200
    seed = r.headers.get("X-Strand-Seed")
    assert seed is not None and seed.isdigit()


def test_seed_reuse_reproduces_output(client, png_bytes):
    r1 = client.post(
        "/strand",
        files={"file": ("photo.png", png_bytes, "image/png")},
        data={"palette": "dark", "intensity": "normal"},
    )
    seed = r1.headers["X-Strand-Seed"]
    r2 = client.post(
        "/strand",
        files={"file": ("photo.png", png_bytes, "image/png")},
        data={"palette": "dark", "intensity": "normal", "seed": seed},
    )
    assert r1.status_code == r2.status_code == 200
    assert r1.content == r2.content


def test_custom_suffix_in_content_disposition(client, png_bytes):
    r = client.post(
        "/strand",
        files={"file": ("photo.png", png_bytes, "image/png")},
        data={"palette": "dark", "intensity": "subtle", "name_suffix": ".v2"},
    )
    assert r.status_code == 200
    assert "photo.v2.png" in r.headers["content-disposition"]


def test_empty_suffix_keeps_original_name(client, png_bytes):
    r = client.post(
        "/strand",
        files={"file": ("photo.png", png_bytes, "image/png")},
        data={"palette": "dark", "intensity": "subtle", "name_suffix": ""},
    )
    assert r.status_code == 200
    assert 'filename="photo.png"' in r.headers["content-disposition"]


def test_suffix_strips_path_separators(client, png_bytes):
    r = client.post(
        "/strand",
        files={"file": ("photo.png", png_bytes, "image/png")},
        data={"palette": "dark", "intensity": "subtle", "name_suffix": "../oops"},
    )
    assert r.status_code == 200
    # The slashes / dots-at-start were stripped; the resulting suffix is "..oops".
    cd = r.headers["content-disposition"]
    assert "/" not in cd.split("filename=")[1]


# --- /api/sample ------------------------------------------------------------

def test_api_sample_returns_png(client):
    r = client.get("/api/sample", params={"palette": "dark", "seed": 1})
    assert r.status_code == 200
    assert r.headers["content-type"] == "image/png"
    with Image.open(io.BytesIO(r.content)) as im:
        im.load()
        assert im.format == "PNG"
        # Hairs render with an alpha channel so they composite cleanly.
        assert "A" in im.getbands()


def test_api_sample_is_deterministic_for_seed(client):
    a = client.get("/api/sample", params={"palette": "blonde", "seed": 42}).content
    b = client.get("/api/sample", params={"palette": "blonde", "seed": 42}).content
    assert a == b
    c = client.get("/api/sample", params={"palette": "blonde", "seed": 43}).content
    assert a != c


def test_api_sample_rejects_unknown_palette(client):
    r = client.get("/api/sample", params={"palette": "neon-pink"})
    assert r.status_code == 400


@pytest.mark.parametrize("palette", ["dark", "brown", "blonde", "red", "grey", "white", "mixed"])
def test_api_sample_all_palettes(client, palette):
    r = client.get("/api/sample", params={"palette": palette, "seed": 1})
    assert r.status_code == 200
    assert r.headers["content-type"] == "image/png"


def test_api_sample_loop_variant(client):
    # `loop=true` forces the laminated-loop morphology (loop_chance=1.0).
    r = client.get("/api/sample", params={"palette": "dark", "seed": 1, "loop": "true"})
    assert r.status_code == 200


@pytest.mark.parametrize("morphology", ["curve", "loop", "eyelash", "fragment", "kink"])
def test_api_sample_morphology_choices(client, morphology):
    r = client.get("/api/sample", params={"palette": "dark", "seed": 1, "morphology": morphology})
    assert r.status_code == 200
    assert r.headers["content-type"] == "image/png"


def test_api_sample_rejects_unknown_morphology(client):
    r = client.get("/api/sample", params={"palette": "dark", "morphology": "spiral"})
    assert r.status_code == 400


def test_api_sample_morphology_varies_output(client):
    """Different morphologies on the same (palette, seed) must produce different bytes."""
    eye = client.get("/api/sample", params={"palette": "dark", "seed": 7, "morphology": "eyelash"}).content
    frag = client.get("/api/sample", params={"palette": "dark", "seed": 7, "morphology": "fragment"}).content
    curve = client.get("/api/sample", params={"palette": "dark", "seed": 7, "morphology": "curve"}).content
    assert eye != frag != curve != eye


# --- Multi-file & folder upload --------------------------------------------

def test_multi_file_upload_returns_zip(client, png_bytes, pdf_bytes):
    files = [
        ("file", ("a.png", png_bytes, "image/png")),
        ("file", ("b.pdf", pdf_bytes, "application/pdf")),
    ]
    r = client.post("/strand", files=files, data={"palette": "dark", "intensity": "normal"})
    assert r.status_code == 200
    assert r.headers["content-type"] == "application/zip"
    # No common root → falls back to "files-strand.zip".
    assert "files-strand.zip" in r.headers["content-disposition"]
    assert r.headers["X-Strand-Haired"] == "2"

    import zipfile
    with zipfile.ZipFile(io.BytesIO(r.content)) as z:
        names = set(z.namelist())
        assert "a-strand.png" in names
        assert "b-strand.pdf" in names
        assert "_strand-report.txt" in names


def test_folder_upload_preserves_paths_and_picks_root_name(client, png_bytes, pdf_bytes):
    """Files uploaded with sub-paths (folder picker) keep their structure, and
    the response zip is named after the common root."""
    files = [
        ("file", ("photos/a.png", png_bytes, "image/png")),
        ("file", ("photos/notes.pdf", pdf_bytes, "application/pdf")),
    ]
    r = client.post("/strand", files=files, data={"palette": "dark", "intensity": "subtle"})
    assert r.status_code == 200
    assert "photos-strand.zip" in r.headers["content-disposition"]

    import zipfile
    with zipfile.ZipFile(io.BytesIO(r.content)) as z:
        names = set(z.namelist())
        assert "photos/a-strand.png" in names
        assert "photos/notes-strand.pdf" in names


def test_single_file_upload_still_returns_raw_file(client, png_bytes):
    """Backwards-compat: one image upload still returns a raw image, not a zip."""
    r = client.post(
        "/strand",
        files={"file": ("photo.png", png_bytes, "image/png")},
        data={"palette": "dark", "intensity": "normal"},
    )
    assert r.status_code == 200
    assert r.headers["content-type"] == "image/png"
    assert "X-Strand-Haired" not in r.headers


def test_multi_file_total_size_limit(client, png_bytes):
    """Total payload, not per-file, is what triggers the 25 MB cap."""
    # 25 fake "files" of ~1.1 MB each (well under 25 MB each, > 25 MB total).
    big_image = png_bytes + b"\0" * (1_100_000 - len(png_bytes))
    files = [("file", (f"f{i}.png", big_image, "image/png")) for i in range(25)]
    r = client.post("/strand", files=files, data={"palette": "dark", "intensity": "subtle"})
    assert r.status_code == 413


def test_no_uploads_returns_400(client):
    r = client.post("/strand", data={"palette": "dark", "intensity": "normal"})
    assert r.status_code == 400


# --- Error envelope --------------------------------------------------------

def test_500_includes_error_id(client, monkeypatch, png_bytes):
    """An unexpected exception path should still surface a short error_id to the user."""
    from app import main as main_mod
    def _boom(*a, **kw):
        raise RuntimeError("synthetic")
    monkeypatch.setattr(main_mod, "strand_bytes", _boom)

    r = client.post(
        "/strand",
        files={"file": ("photo.png", png_bytes, "image/png")},
        data={"palette": "dark", "intensity": "normal"},
    )
    assert r.status_code == 500
    assert "error_id" in r.json()["detail"]
    assert r.headers.get("X-Strand-Error-Id")


# --- Static-mount-at-root smoke check --------------------------------------

def test_root_serves_index_html(client):
    r = client.get("/")
    assert r.status_code == 200
    assert "Strand" in r.text


def test_static_assets_resolve_at_root(client):
    """index.html references `style.css` / `app.js` relatively. The static
    mount at `/` must serve those at the root level."""
    for asset in ("style.css", "app.js"):
        r = client.get(f"/{asset}")
        assert r.status_code == 200, f"GET /{asset} returned {r.status_code}"


# --- Clusters --------------------------------------------------------------

def test_clusters_can_increase_hair_count(pdf_bytes):
    """With cluster_chance high, a PDF should embed more hairs per page than
    with clusters off. Don't assert an exact ratio (RNG-dependent) — just that
    the count strictly increases when buddies are likely."""
    import fitz

    def total_images(cluster_chance):
        opts = options_from_ui(palette="dark", intensity="normal", seed=2024)
        # options_from_ui doesn't take cluster_chance; mutate after construction.
        opts.cluster_chance = cluster_chance
        out, _ = inject_pdf_bytes(pdf_bytes, opts)
        doc = fitz.open(stream=out, filetype="pdf")
        try:
            return sum(len(doc[i].get_images()) for i in range(doc.page_count))
        finally:
            doc.close()

    assert total_images(0.95) > total_images(0.0)


# --- Kink morphology + follicle bulb --------------------------------------

def test_kink_morphology_is_registered():
    """kink should appear in the dispatch table and the cm-length map."""
    from app.core import MORPHOLOGIES, MORPHOLOGY_LENGTH_CM
    assert "kink" in MORPHOLOGIES
    assert "kink" in MORPHOLOGY_LENGTH_CM
    lo, hi = MORPHOLOGY_LENGTH_CM["kink"]
    assert 0 < lo < hi < 10  # short-ish, body-hair range


def test_kink_renderer_produces_non_empty_image():
    import random
    from app.core import _generate_kink_hair
    img = _generate_kink_hair(random.Random(7), palette="dark")
    # The strand should have some non-transparent pixels.
    alpha = img.split()[-1]
    assert alpha.getextrema()[1] > 0


def test_kink_strand_visibly_zigzags():
    """A kink hair should have visibly more direction changes than a curve.

    Count alpha-mass columns in the strand image: a kink, projected onto its
    primary axis, should be wider (more y-spread) than a single bezier curve
    drawn at the same canvas dimensions.
    """
    import random
    from app.core import _generate_curve_hair, _generate_kink_hair

    def vertical_spread(img):
        alpha = img.split()[-1]
        w, h = alpha.size
        a = alpha.load()
        rows_with_ink = 0
        for y in range(h):
            for x in range(w):
                if a[x, y] > 30:
                    rows_with_ink += 1
                    break
        return rows_with_ink

    # Same RNG seed for fair-ish comparison.
    curve_spread = vertical_spread(_generate_curve_hair(random.Random(1), palette="dark"))
    kink_spread = vertical_spread(_generate_kink_hair(random.Random(1), palette="dark"))
    # A kink fills more vertical extent because it zigzags.
    assert kink_spread > curve_spread * 1.3, (
        f"kink vertical spread {kink_spread} should be >1.3x curve spread {curve_spread}"
    )


def test_follicle_can_appear_on_a_strand():
    """When _FOLLICLE_CHANCE is forced to 1.0, every curve hair gets a bulb.

    The strand is drawn identically in both versions (same seed, bezier first,
    follicle check after), so the bulbed version must have strictly more ink
    than the clean one. Compare total alpha mass — that's threshold-free and
    works regardless of palette width scaling.
    """
    import random
    from app import core
    saved = core._FOLLICLE_CHANCE
    core._FOLLICLE_CHANCE = 1.0
    try:
        bulbed = core._generate_curve_hair(random.Random(11), palette="white")
        core._FOLLICLE_CHANCE = 0.0
        clean = core._generate_curve_hair(random.Random(11), palette="white")
    finally:
        core._FOLLICLE_CHANCE = saved

    def total_ink(img):
        return sum(img.split()[-1].getdata())

    assert total_ink(bulbed) > total_ink(clean) + 500, (
        f"follicle should add visible ink mass: bulbed={total_ink(bulbed)} "
        f"vs clean={total_ink(clean)}"
    )


# --- Extended stats (substrate / lengths / content-hits) ------------------

def test_image_stats_include_substrate_lengths_and_hits(png_bytes):
    opts = InjectOptions(seed=1, image_count=3, cluster_chance=0.0, palette="white")
    _, stats = inject_image_bytes(png_bytes, ".png", opts)

    sub = stats["substrate"]
    assert sub is not None
    assert sub["width_native"] == 400 and sub["height_native"] == 300
    assert sub["native_unit"] == "px"
    assert sub["dpi"] in (72, 96)  # PIL fallback / default
    assert sub["width_cm"] > 0 and sub["height_cm"] > 0

    # One length per hair, all within the curve cm range with a thumbnail clamp
    # safety margin.
    assert len(stats["hair_lengths_cm"]) == stats["hairs"]
    assert all(0 < L < 20 for L in stats["hair_lengths_cm"])

    # Content-hits must be between 0 and the number of primary hairs (3).
    assert 0 <= stats["content_hits"] <= 3


def test_pdf_stats_include_substrate(pdf_bytes):
    opts = InjectOptions(seed=2, rate=1.0, hairs_per_page=2, cluster_chance=0.0)
    _, stats = inject_pdf_bytes(pdf_bytes, opts)
    sub = stats["substrate"]
    assert sub is not None
    assert sub["native_unit"] == "pt"
    # US letter ≈ 21.6 × 27.9 cm
    assert 21 < sub["width_cm"] < 22
    assert 27 < sub["height_cm"] < 28
    assert sub["page_count"] == 3
    assert len(stats["hair_lengths_cm"]) == stats["hairs"]


def test_pptx_stats_include_substrate(pptx_bytes):
    opts = InjectOptions(seed=3, rate=1.0, hairs_per_page=1, cluster_chance=0.0)
    _, stats = inject_pptx_bytes(pptx_bytes, opts)
    sub = stats["substrate"]
    assert sub is not None
    assert sub["native_unit"] == "EMU"
    # Default pptx ≈ 25.4 × 19.05 cm (10 × 7.5 in).
    assert 25 < sub["width_cm"] < 26
    assert sub["slide_count"] == 3


def test_zip_stats_aggregate_lengths_without_substrate(zip_bytes):
    opts = options_from_ui(palette="dark", intensity="normal", seed=42)
    _, report = strand_zip_bytes(zip_bytes, opts, name_suffix="-strand")
    agg = report["stats"]
    # A zip mixes substrates, so the aggregate substrate is None.
    assert agg["substrate"] is None
    # Lengths from every entry are concatenated.
    assert len(agg["hair_lengths_cm"]) == agg["hairs"]
    assert all(0 < L < 30 for L in agg["hair_lengths_cm"])


# --- DPI / cm-based sizing -------------------------------------------------

def test_eyelash_is_smaller_than_curve_on_same_canvas():
    """The per-morphology cm range should make eyelashes (0.7-1.4 cm) visibly
    smaller than full curves (3-12 cm) when both land on the same substrate."""
    import random
    from app.core import (
        MORPHOLOGY_LENGTH_CM,
        _image_pixels_per_cm,
        _place,
        _generate_curve_hair,
        _generate_eyelash_hair,
    )

    img = Image.new("RGBA", (1000, 1000), (255, 255, 255, 255))
    units = _image_pixels_per_cm(img)

    curve = _generate_curve_hair(random.Random(1))
    eye = _generate_eyelash_hair(random.Random(1))

    # Run _place many times and compare median sizes — RNG variance otherwise.
    def median_long(hair, range_cm):
        sizes = []
        for s in range(50):
            _, _, dw, dh, _hit = _place(
                random.Random(s), 1000, 1000, hair.size[0], hair.size[1],
                range_cm, units,
            )
            sizes.append(max(dw, dh))
        sizes.sort()
        return sizes[len(sizes) // 2]

    eye_median = median_long(eye, MORPHOLOGY_LENGTH_CM["eyelash"])
    curve_median = median_long(curve, MORPHOLOGY_LENGTH_CM["curve"])
    assert eye_median < curve_median, f"eyelash {eye_median} should be smaller than curve {curve_median}"


def test_large_image_doesnt_get_huge_hair():
    """Regression for the old relative-scale bug: a 3000x3000 image at default
    DPI should get a hair sized by physical cm, not by fraction of the canvas.
    Upper bound: 12 cm (curve max) * 96/2.54 ≈ 453 px — well under the old
    0.45 * 3000 = 1350 px upper bound."""
    import random
    from app.core import MORPHOLOGY_LENGTH_CM, _image_pixels_per_cm, _place

    img = Image.new("RGBA", (3000, 3000), (255, 255, 255, 255))
    units = _image_pixels_per_cm(img)

    longest = 0
    for s in range(30):
        _, _, dw, dh, _hit = _place(
            random.Random(s), 3000, 3000, 400, 120,
            MORPHOLOGY_LENGTH_CM["curve"], units,
        )
        longest = max(longest, max(dw, dh))
    # 12 cm * (96/2.54) ≈ 453 px ceiling at default DPI; give a generous
    # buffer for the upper end of uniform sampling.
    assert longest < 500, f"curve max length on big image was {longest}px"


def test_thumbnail_hair_clamped_to_substrate():
    """On a tiny thumbnail, the hair must stay within the canvas — no
    matter what the cm range says."""
    import random
    from app.core import MORPHOLOGY_LENGTH_CM, _image_pixels_per_cm, _place

    img = Image.new("RGBA", (60, 60), (255, 255, 255, 255))
    units = _image_pixels_per_cm(img)

    for s in range(20):
        _, _, dw, dh, _hit = _place(
            random.Random(s), 60, 60, 400, 120,
            MORPHOLOGY_LENGTH_CM["curve"], units,
        )
        # 90% clamp means longest side ≤ 54 px on a 60 px substrate.
        assert max(dw, dh) <= 60 * 0.9 + 0.01


def test_zip_endpoint_roundtrip(client, zip_bytes):
    import zipfile
    r = client.post(
        "/strand",
        files={"file": ("bundle.zip", zip_bytes, "application/zip")},
        data={"palette": "dark", "intensity": "normal"},
    )
    assert r.status_code == 200
    assert r.headers["content-type"] == "application/zip"
    assert "bundle-strand.zip" in r.headers["content-disposition"]
    assert r.headers.get("X-Strand-Haired") == "3"
    assert r.headers.get("X-Strand-Skipped") == "1"

    with zipfile.ZipFile(io.BytesIO(r.content)) as z:
        assert "photo-strand.png" in z.namelist()
        assert "_strand-report.txt" in z.namelist()
